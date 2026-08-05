"""Figures out of slide decks.

A deck keeps its meaning in pictures. Parsed as text alone, a forty-slide
architecture review yields forty titles and nothing worth retrieving — which is
the failure this path exists to fix, and it is a silent one: the deck ingests
cleanly and simply has no content.

Two things are load-bearing here and both are easy to break without noticing.

Extraction order must match the order docling writes `<!-- image -->` markers,
because the two are paired positionally. Get it wrong and every figure's
description lands on a different figure's position — with no error, and no way to
see it except by reading the output.

And template furniture must stay out. A logo on the slide master is not a figure;
described once per slide it is forty model calls and forty copies of the same
sentence in the retrieval index.
"""

from __future__ import annotations

import logging

import pytest

pytest.importorskip("pptx", reason="python-pptx arrives with the docling extra")

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from prismdoc.ooxml import extract_pptx_images  # noqa: E402
from prismdoc.stages.figures import _replace_image_markers  # noqa: E402


def png(tmp_path, name: str, size: tuple[int, int] = (200, 150)):
    from PIL import Image

    path = tmp_path / name
    Image.new("RGB", size, (200, 40, 40)).save(path)
    return str(path)


def deck(tmp_path, slides: list[list[str]], name: str = "deck.pptx") -> str:
    """Build a deck: one entry per slide, listing the images placed on it."""
    presentation = Presentation()
    for index, images in enumerate(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Slide {index}"
        for offset, image in enumerate(images):
            slide.shapes.add_picture(
                image, Inches(1 + offset * 3), Inches(2), Inches(2)
            )
    path = tmp_path / name
    presentation.save(path)
    return str(path)


class TestExtractionOrder:
    """The contract `_replace_image_markers` pairs against."""

    def test_pictures_come_back_in_reading_order(self, tmp_path):
        first, second, third = (
            png(tmp_path, "a.png", (200, 150)),
            png(tmp_path, "b.png", (240, 150)),
            png(tmp_path, "c.png", (280, 150)),
        )
        path = deck(tmp_path, [[first, second], [], [third]])

        images = extract_pptx_images(path)

        assert [(i.page_index, i.width) for i in images] == [
            (0, 200),
            (0, 240),
            (2, 280),
        ]

    def test_a_slide_without_pictures_contributes_none(self, tmp_path):
        path = deck(tmp_path, [[], [png(tmp_path, "a.png")], []])

        assert [i.page_index for i in extract_pptx_images(path)] == [1]

    def test_pictures_inside_a_group_keep_their_place(self, tmp_path):
        """Docling walks into groups too, so skipping them would shift the pairing."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        path = deck(tmp_path, [[png(tmp_path, "a.png"), png(tmp_path, "b.png")]])
        presentation = Presentation(path)
        slide = presentation.slides[0]
        pictures = [
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        assert len(pictures) == 2, "fixture did not produce two pictures"

        assert len(extract_pptx_images(path)) == 2


class TestWhatIsNotAFigure:
    def test_template_furniture_is_excluded(self, tmp_path):
        """Only pictures placed on slides — not on layouts or the master."""
        path = deck(tmp_path, [[], [], []])

        assert extract_pptx_images(path) == []

    def test_a_decorative_sliver_is_skipped(self, tmp_path):
        """Bullet glyphs and rules cost a model call and return noise."""
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "One"
        slide.shapes.add_picture(
            png(tmp_path, "rule.png", (400, 4)), Inches(1), Inches(2),
            width=Inches(4), height=Pt(3),
        )
        path = tmp_path / "thin.pptx"
        presentation.save(path)

        assert extract_pptx_images(str(path)) == []


class TestGeometry:
    def test_bbox_is_in_points_from_the_top_left(self, tmp_path):
        """Same convention as prismdoc.pdf, so Figure.bbox means one thing."""
        path = deck(tmp_path, [[png(tmp_path, "a.png")]])

        bbox = extract_pptx_images(path)[0].bbox

        assert bbox is not None
        # Placed at 1in from the left, 2in from the top → 72pt, 144pt.
        assert bbox[0] == pytest.approx(72, abs=1)
        assert bbox[1] == pytest.approx(144, abs=1)
        assert bbox[2] > bbox[0] and bbox[3] > bbox[1]

    def test_dimensions_are_the_stored_pixels(self, tmp_path):
        """No rasterising happens, so the bytes are the original asset."""
        path = deck(tmp_path, [[png(tmp_path, "a.png", (321, 123))]])

        image = extract_pptx_images(path)[0]

        assert (image.width, image.height) == (321, 123)
        assert image.mime == "image/png"


class TestMarkerPlacement:
    MARKDOWN = "# One\n\n<!-- image -->\n\n<!-- image -->\n\n# Two\n\n<!-- image -->"

    def test_tokens_land_where_the_pictures_were(self):
        result = _replace_image_markers(self.MARKDOWN, ["a", "b", "c"])

        assert "<!-- image -->" not in result
        assert result.index("[[FIGURE:a]]") < result.index("[[FIGURE:b]]")
        assert result.index("[[FIGURE:b]]") < result.index("# Two")
        assert result.index("[[FIGURE:c]]") > result.index("# Two")

    def test_surplus_markers_are_left_alone(self, caplog):
        """Docling marks pictures the extractor skipped; an invented token would
        merge as '[unprocessed figure ...]' and read like a bug in the corpus."""
        with caplog.at_level(logging.WARNING, logger="prismdoc.stages.figures"):
            result = _replace_image_markers(self.MARKDOWN, ["a"])

        assert result.count("<!-- image -->") == 2
        assert "[[FIGURE:a]]" in result
        assert "3 image marker(s) but 1 extracted" in caplog.text

    def test_surplus_figures_are_appended_rather_than_dropped(self, caplog):
        with caplog.at_level(logging.WARNING, logger="prismdoc.stages.figures"):
            result = _replace_image_markers("# One\n\n<!-- image -->", ["a", "b"])

        assert "[[FIGURE:a]]" in result
        assert result.rstrip().endswith("[[FIGURE:b]]")
        assert "1 image marker(s) but 2 extracted" in caplog.text

    def test_no_markers_at_all_says_so(self, caplog):
        with caplog.at_level(logging.WARNING, logger="prismdoc.stages.figures"):
            result = _replace_image_markers("Flat text.", ["a"])

        assert "[[FIGURE:a]]" in result
        assert "No image markers" in caplog.text

    def test_markdown_is_untouched_when_there_are_no_figures(self):
        assert _replace_image_markers(self.MARKDOWN, []) == self.MARKDOWN
